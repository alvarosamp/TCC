"""Gera slides do TCC em formato .pptx — versao final com diagramas."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── Paleta ────────────────────────────────────────────────────────────────────
AZ  = RGBColor(0x1A, 0x37, 0x5E)   # azul escuro
AM  = RGBColor(0x1F, 0x6F, 0xAB)   # azul medio
AC  = RGBColor(0x5B, 0xA3, 0xD0)   # azul claro
VE  = RGBColor(0x1E, 0x8B, 0x4C)   # verde
VA  = RGBColor(0xD4, 0xEF, 0xDF)   # verde agua claro
VR  = RGBColor(0xC0, 0x39, 0x2B)   # vermelho
LA  = RGBColor(0xFD, 0xF2, 0xF8)   # lavanda muito claro
AM_ = RGBColor(0xF3, 0x9C, 0x12)   # amarelo/laranja
CZ  = RGBColor(0xEC, 0xF0, 0xF1)   # cinza claro
BR  = RGBColor(0xFF, 0xFF, 0xFF)   # branco
PT  = RGBColor(0x1C, 0x1C, 0x1C)   # preto suave
CZ2 = RGBColor(0x7F, 0x8C, 0x8D)   # cinza medio

W = Inches(13.33)
H = Inches(7.5)

# ── Primitivas ────────────────────────────────────────────────────────────────

def slide_branco(prs, cor_bg=BR):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = cor_bg
    return s

def box(slide, x, y, w, h, fill=None, border=None, border_w=1, radius=0):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(border_w)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, size=14, bold=False, color=PT,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf.text_frame.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tf

def header(slide, titulo, subtitulo="", cor=AZ):
    box(slide, 0, 0, 13.33, 1.2, fill=cor)
    box(slide, 0, 1.2, 13.33, 0.05, fill=AM_)
    txt(slide, titulo, 0.35, 0.1, 12.5, 0.7, size=26, bold=True, color=BR)
    if subtitulo:
        txt(slide, subtitulo, 0.35, 0.72, 12.5, 0.42, size=13,
            color=RGBColor(0xB8, 0xD4, 0xEA))

def numero_slide(slide, n, total):
    txt(slide, f"{n} / {total}", 12.5, 7.2, 0.8, 0.3, size=9,
        color=CZ2, align=PP_ALIGN.RIGHT)

def chip(slide, label, x, y, w=1.5, h=0.45, fill=AM, cor_txt=BR, size=11, bold=True):
    box(slide, x, y, w, h, fill=fill)
    txt(slide, label, x, y, w, h, size=size, bold=bold, color=cor_txt,
        align=PP_ALIGN.CENTER)

def seta(slide, x, y, horizontal=True, cor=AM):
    if horizontal:
        box(slide, x, y+0.16, 0.35, 0.1, fill=cor)
        box(slide, x+0.22, y+0.05, 0.13, 0.32, fill=cor)
    else:
        box(slide, x+0.13, y, 0.1, 0.35, fill=cor)
        box(slide, x, y+0.22, 0.36, 0.13, fill=cor)

def linha_tabela(slide, cols, valores, y, alts, cor_bg, cor_txt=PT, size=11,
                 bold=False, bordas=True):
    x = 0.3
    for i, (v, w) in enumerate(zip(valores, alts)):
        box(slide, x, y, w-0.03, 0.52,
            fill=cor_bg,
            border=RGBColor(0xCC,0xCC,0xCC) if bordas else None,
            border_w=0.5)
        al = PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT
        txt(slide, str(v), x+0.05, y+0.06, w-0.1, 0.42,
            size=size, bold=bold, color=cor_txt, align=al)
        x += w

# ── Slides ────────────────────────────────────────────────────────────────────

def slide_capa(prs):
    s = slide_branco(prs, cor_bg=AZ)
    # faixa decorativa lateral
    box(s, 0, 0, 0.18, 7.5, fill=AM_)
    # gradiente simulado com caixa clara
    box(s, 0.18, 0, 13.15, 7.5, fill=AZ)
    box(s, 0.18, 5.3, 13.15, 2.2, fill=AM)

    txt(s, "TCC — TinyML / MLOps", 0.5, 0.6, 12.3, 0.8,
        size=16, color=AM_, bold=True)
    txt(s,
        "Pipeline Generico para Deteccao de\nAnomalias em Series Temporais",
        0.5, 1.35, 12.3, 2.0, size=38, bold=True, color=BR,
        align=PP_ALIGN.LEFT)
    txt(s,
        "Estudo de caso: dados sismicos  ·  Validacao em ESP32 real  ·  OTA via HTTP",
        0.5, 3.35, 12.3, 0.6, size=15,
        color=RGBColor(0xA8,0xC8,0xE8))

    # badges
    badges = [("AUC-PR 0.877", VE), ("15.377 params", AM),
              ("ESP32-D0WD-V3", AM_), ("OTA via HTTP", VR)]
    for i,(b,c) in enumerate(badges):
        chip(s, b, 0.5 + i*2.9, 4.15, w=2.6, h=0.5, fill=c, size=12)

    txt(s, "Engenharia / Ciencia de Dados  ·  2026",
        0.5, 5.55, 12.3, 0.5, size=14, color=BR)
    txt(s, "github.com/alvarosamp/TCC",
        0.5, 6.05, 12.3, 0.4, size=12,
        color=RGBColor(0x90,0xB8,0xD8))


def slide_problema(prs, n, tot):
    s = slide_branco(prs, CZ)
    header(s, "O Problema", "Por que processar na borda?")
    numero_slide(s, n, tot)

    # dois lados: sem borda vs com borda
    colunas = [
        ("Sem TinyML na borda", VR,
         ["Sensor envia 100% dos dados",
          "Transmissao continua (energia/custo)",
          "Latencia: dados chegam ao servidor",
          "Infraestrutura cara para escalar",
          "Privacidade: dados brutos expostos"]),
        ("Com TinyML na borda", VE,
         ["ESP32 decide localmente (20ms/janela)",
          "Transmite apenas quando ha anomalia",
          "Latencia zero — decisao embarcada",
          "Escalavel: logica no dispositivo",
          "Apenas eventos chegam ao servidor"]),
    ]
    for col,(titulo,cor,items) in enumerate(colunas):
        x = 0.35 + col*6.6
        box(s, x, 1.4, 6.3, 5.6, fill=BR,
            border=cor, border_w=2)
        box(s, x, 1.4, 6.3, 0.65, fill=cor)
        txt(s, titulo, x+0.1, 1.43, 6.1, 0.58,
            size=15, bold=True, color=BR, align=PP_ALIGN.CENTER)
        for i,item in enumerate(items):
            icon = "✗" if col==0 else "✓"
            cor_i = VR if col==0 else VE
            txt(s, icon, x+0.2, 2.25+i*0.9, 0.3, 0.5,
                size=18, bold=True, color=cor_i)
            txt(s, item, x+0.55, 2.25+i*0.9, 5.6, 0.5,
                size=13, color=PT)

    txt(s, "Reducao de transmissao, armazenamento, energia e custo operacional",
        1.5, 7.1, 10.3, 0.35, size=12, italic=True,
        color=AM, align=PP_ALIGN.CENTER)


def slide_arquitetura(prs, n, tot):
    s = slide_branco(prs, CZ)
    header(s, "Arquitetura do Sistema", "Pipeline de ponta a ponta")
    numero_slide(s, n, tot)

    # tres zonas
    zonas = [
        ("DADOS & FEATURES", AM, 0.25, 1.4, 3.8, 5.7),
        ("MLOPS & MODELOS",  AZ, 4.4,  1.4, 4.6, 5.7),
        ("BORDA & OTA",      VE, 9.3,  1.4, 3.8, 5.7),
    ]
    for titulo, cor, x, y, w, h in zonas:
        box(s, x, y, w, h, fill=BR, border=cor, border_w=1)
        box(s, x, y, w, 0.42, fill=cor)
        txt(s, titulo, x+0.1, y+0.05, w-0.2, 0.35,
            size=10, bold=True, color=BR, align=PP_ALIGN.CENTER)

    # blocos zona 1
    b1 = [("MiniSEED\n(raw)", 1.2), ("Adapter\nDominio", 1.2),
          ("Dataset NPZ\n(contrato)", 1.2), ("Features\nTabulares", 1.2)]
    for i,(label,_) in enumerate(b1):
        chip(s, label, 0.35, 2.05+i*1.2, w=3.6, h=0.85,
             fill=AM, size=10)
        if i < len(b1)-1:
            seta(s, 1.75, 2.92+i*1.2, horizontal=False)

    # blocos zona 2
    b2 = [("Treino\n(MLflow)", AM_), ("HPO\n(Optuna)", AM_),
          ("Quality\nGate", VE), ("Exportacao\nTFLite", AC)]
    for i,(label,cor) in enumerate(b2):
        chip(s, label, 4.5, 2.05+i*1.2, w=4.2, h=0.85,
             fill=cor, size=10)
        if i < len(b2)-1:
            seta(s, 6.3, 2.92+i*1.2, horizontal=False)

    # blocos zona 3
    b3 = [("Servidor\nFastAPI", AM), ("OTA\nHTTP", AM_),
          ("SPIFFS\n(flash)", AC), ("TFLite Micro\n(inferencia)", VE)]
    for i,(label,cor) in enumerate(b3):
        chip(s, label, 9.4, 2.05+i*1.2, w=3.6, h=0.85,
             fill=cor, size=10)
        if i < len(b3)-1:
            seta(s, 11.0, 2.92+i*1.2, horizontal=False)

    # setas entre zonas
    for x in [4.08, 9.13]:
        box(s, x, 3.7, 0.3, 0.08, fill=AM_)
        box(s, x+0.18, 3.6, 0.12, 0.28, fill=AM_)


def slide_dataset(prs, n, tot):
    s = slide_branco(prs, CZ)
    header(s, "Dataset Sismico", "MiniSEED — janelamento e preprocessamento edge-aware")
    numero_slide(s, n, tot)

    # info cards
    cards = [
        ("Origem", "Estacoes sismologicas\n(formato MiniSEED)", AM),
        ("Classes", "Normal (background)\nAnomalo (evento sismico)", AZ),
        ("Taxa", "40 Hz\n(800 amostras = 20 s)", AM_),
        ("Overlap", "50%\n(step de 10 s)", VE),
        ("Split", "Por evento\n(sem vazamento temporal)", VR),
    ]
    for i,(k,v,c) in enumerate(cards):
        x = 0.3 + i*2.6
        box(s, x, 1.45, 2.4, 1.5, fill=BR, border=c, border_w=2)
        box(s, x, 1.45, 2.4, 0.45, fill=c)
        txt(s, k, x+0.08, 1.47, 2.25, 0.4,
            size=11, bold=True, color=BR, align=PP_ALIGN.CENTER)
        txt(s, v, x+0.08, 1.97, 2.25, 0.9,
            size=11, color=PT, align=PP_ALIGN.CENTER)

    # diagrama de janelamento
    box(s, 0.3, 3.15, 12.7, 1.6, fill=BR, border=AM, border_w=1)
    txt(s, "Janelamento com overlap 50%", 0.4, 3.18, 5.0, 0.35,
        size=11, bold=True, color=AM)

    # barra do sinal
    box(s, 0.4, 3.6, 12.4, 0.35, fill=CZ, border=AM, border_w=1)
    txt(s, "Sinal continuo (horas / dias)", 5.5, 3.62, 4.0, 0.3,
        size=10, color=CZ2, align=PP_ALIGN.CENTER)

    # janelas
    janelas = [(0.4,3.0,"J0\n20s"), (1.4,3.0,"J1\n20s"),
               (2.4,3.0,"J2\n20s"), (3.4,3.0,"J3\n20s")]
    cores_j = [AM, AC, AM, AC]
    for i,((ox,_,lbl),c) in enumerate(zip(janelas,cores_j)):
        bx = 0.4 + i*1.25
        box(s, bx, 4.05, 2.5, 0.52, fill=c, border=AZ, border_w=1)
        txt(s, lbl, bx+0.05, 4.1, 2.4, 0.42,
            size=10, bold=True, color=BR, align=PP_ALIGN.CENTER)
    txt(s, "← step 10s →", 1.5, 4.6, 1.5, 0.3, size=9, color=CZ2)
    txt(s, "← step 10s →", 2.75, 4.6, 1.5, 0.3, size=9, color=CZ2)

    # preprocessamento
    box(s, 0.3, 4.95, 12.7, 2.3, fill=BR, border=AZ, border_w=1)
    txt(s, "Preprocessamento Edge-Aware (reproduzivel no microcontrolador)",
        0.4, 4.97, 10.0, 0.38, size=11, bold=True, color=AZ)

    steps = ["resample\n40 Hz", "detrend\nlinear", "demean",
             "taper\n5%", "bandpass\n0.5-15 Hz", "zscore\npor janela"]
    cores_s = [AM, AM, AC, AM, AZ, VE]
    for i,(st,c) in enumerate(zip(steps, cores_s)):
        x = 0.4 + i*2.12
        chip(s, st, x, 5.42, w=1.95, h=0.75, fill=c, size=9)
        if i < len(steps)-1:
            txt(s, "→", x+1.97, 5.6, 0.2, 0.3,
                size=14, bold=True, color=AM, align=PP_ALIGN.CENTER)
    txt(s, "remove_response EXCLUIDO: depende de StationXML — inviavel em firmware",
        0.4, 6.25, 12.5, 0.35, size=10, italic=True, color=CZ2)


def slide_modelos_arch(prs, n, tot):
    s = slide_branco(prs, CZ)
    header(s, "Arquitetura dos Modelos", "Familias suportadas e candidato edge")
    numero_slide(s, n, tot)

    familias = [
        ("Classicos\nSupervisionados", AM,
         "Random Forest\nExtra Trees\nLogistic Regression",
         "Baselines e interpretabilidade"),
        ("Classicos Nao\nSupervisionados", AZ,
         "Isolation Forest",
         "Cenarios com poucos labels"),
        ("Redes Leves\n(TinyML)", VE,
         "Tiny CNN  ← candidato atual\nTiny TCN\nLSTM",
         "Modelos para exportacao TFLite"),
        ("Autoencoders", AM_,
         "Dense AE\nCNN AE",
         "Deteccao por erro de reconstrucao"),
    ]
    for i,(titulo,cor,modelos,uso) in enumerate(familias):
        x = 0.3 + i*3.25
        box(s, x, 1.45, 3.05, 3.2, fill=BR, border=cor, border_w=2)
        box(s, x, 1.45, 3.05, 0.55, fill=cor)
        txt(s, titulo, x+0.08, 1.47, 2.9, 0.5,
            size=11, bold=True, color=BR, align=PP_ALIGN.CENTER)
        txt(s, modelos, x+0.12, 2.1, 2.85, 1.3,
            size=11, color=PT)
        box(s, x, 3.5, 3.05, 0.55, fill=CZ, border=cor, border_w=1)
        txt(s, uso, x+0.08, 3.52, 2.9, 0.5,
            size=9, italic=True, color=CZ2, align=PP_ALIGN.CENTER)

    # arquitetura tiny_cnn
    box(s, 0.3, 4.85, 12.7, 2.45, fill=BR, border=VE, border_w=2)
    box(s, 0.3, 4.85, 12.7, 0.45, fill=VE)
    txt(s, "Tiny CNN — Candidato Edge Atual  (15.377 parametros)",
        0.4, 4.87, 12.5, 0.4, size=12, bold=True, color=BR)

    layers = [
        ("Input\n800x1", CZ2), ("Conv1D\nSep x3", AM),
        ("BN+Drop", AC), ("AvgPool\n(global)", VE),
        ("Dense\n48", AM_), ("Sigmoid\nout", VR)
    ]
    for i,(lbl,c) in enumerate(layers):
        x = 0.5 + i*2.12
        chip(s, lbl, x, 5.42, w=1.9, h=0.75, fill=c, size=10)
        if i < len(layers)-1:
            txt(s, "→", x+1.92, 5.6, 0.22, 0.3,
                size=14, bold=True, color=AM, align=PP_ALIGN.CENTER)

    params_txt = "head_pooling: avg  ·  conv_type: conv  ·  n_blocks: 3  ·  base_filters: 16  ·  dropout: 0.24"
    txt(s, params_txt, 0.4, 6.28, 12.5, 0.35,
        size=10, color=CZ2, align=PP_ALIGN.CENTER)


def slide_treinamento(prs, n, tot):
    s = slide_branco(prs, CZ)
    header(s, "Treinamento e HPO", "MLflow para rastreamento · Optuna para busca de hiperparametros")
    numero_slide(s, n, tot)

    # fluxo mlflow
    box(s, 0.3, 1.42, 5.8, 5.85, fill=BR, border=AM, border_w=1)
    txt(s, "Fluxo MLflow", 0.4, 1.45, 5.6, 0.38,
        size=12, bold=True, color=AM)

    mlf = ["Carrega profile\n(seismic_edge_v1)",
           "Carrega dataset NPZ",
           "Extrai features\ntabulares (26 features)",
           "Treina modelos\nenabled=true",
           "Avalia val + teste\n(AUC-PR, F1, FP/h)",
           "Registra no MLflow\n(params + metricas)",
           "Seleciona candidato\nedge_candidate=true",
           "Salva\ncandidate_manifest.json"]
    for i,item in enumerate(mlf):
        cor = [AM,AC,AM,AM_,VE,AM,VE,VR][i]
        chip(s, item, 0.4, 1.95+i*0.65, w=5.5, h=0.6,
             fill=cor, size=9)
        if i < len(mlf)-1:
            txt(s, "↓", 2.95, 2.57+i*0.65, 0.3, 0.2,
                size=10, color=AM, align=PP_ALIGN.CENTER)

    # optuna
    box(s, 6.45, 1.42, 6.55, 2.9, fill=BR, border=AM_, border_w=1)
    txt(s, "HPO com Optuna", 6.55, 1.45, 6.3, 0.38,
        size=12, bold=True, color=AM_)

    opts = [("Trials", "60 rodadas de busca"),
            ("Metrica", "val_auc_pr (maximizar)"),
            ("Espaco", "filters, kernel, lr, dropout..."),
            ("Resultado", "best_value = 0.9671")]
    for i,(k,v) in enumerate(opts):
        box(s, 6.55, 1.97+i*0.6, 1.5, 0.52, fill=AM_)
        txt(s, k, 6.55, 1.97+i*0.6, 1.5, 0.52,
            size=10, bold=True, color=BR, align=PP_ALIGN.CENTER)
        txt(s, v, 8.12, 2.0+i*0.6, 4.7, 0.45, size=10, color=PT)

    # metricas finais
    box(s, 6.45, 4.55, 6.55, 2.72, fill=BR, border=VE, border_w=2)
    box(s, 6.45, 4.55, 6.55, 0.45, fill=VE)
    txt(s, "Resultado Final do Candidato", 6.55, 4.57, 6.3, 0.4,
        size=11, bold=True, color=BR)

    metricas = [("AUC-PR (teste)", "0.8775"),
                ("F1 (teste)", "0.8109"),
                ("Precision", "0.8431"),
                ("Recall", "0.7810"),
                ("FP/h", "6.75"),
                ("Parametros", "15.377")]
    for i,(k,v) in enumerate(metricas):
        y = 5.1 + i*0.35
        txt(s, k, 6.6, y, 4.0, 0.33, size=11, color=PT)
        txt(s, v, 10.5, y, 2.2, 0.33, size=11,
            bold=True, color=VE, align=PP_ALIGN.RIGHT)


def slide_comparativo(prs, n, tot):
    s = slide_branco(prs, CZ)
    header(s, "Comparativo de Modelos",
           "AUC-PR como metrica primaria (deteccao desbalanceada)")
    numero_slide(s, n, tot)

    cols = ["Modelo", "AUC-PR", "F1", "Precision", "Recall", "FP/h", "Edge?"]
    alts = [4.5, 1.3, 1.2, 1.3, 1.2, 1.1, 1.0]

    linha_tabela(s, cols, cols, 1.42, alts, cor_bg=AZ, cor_txt=BR,
                 size=11, bold=True)

    linhas = [
        (["Optuna Tiny CNN v4 (ref. externa)",
          "0.9127","0.8526","0.8982","0.8114","4.90","Sim"], CZ, PT),
        (["Tiny CNN (atual) ← APROVADO",
          "0.8775","0.8109","0.8431","0.7810","6.75","Sim"], VA, VE),
        (["Tiny CNN (baseline sem HPO)",
          "0.8982","0.7951","0.7310","0.8716","16.94","Sim"], BR, PT),
        (["Tiny TCN",
          "0.8964","0.7666","0.6790","0.8801","21.98","Sim"], CZ, PT),
        (["Random Forest (Optuna)",
          "0.8127","0.7367","0.7974","0.6846","9.26","Nao"], BR, PT),
        (["Extra Trees (Optuna)",
          "0.7901","0.7102","0.7589","0.6675","11.30","Nao"], CZ, PT),
        (["STA/LTA (baseline tradicional)",
          "0.1662","0.2760","0.1773","0.6230","—","Nao"], BR, CZ2),
    ]
    for i,(vals,bg,fg) in enumerate(linhas):
        destaque = i==1
        linha_tabela(s, cols, vals, 1.97+i*0.6, alts,
                     cor_bg=bg,
                     cor_txt=VE if destaque else fg,
                     size=10, bold=destaque)

    txt(s, "Quality gate: AUC-PR >= 0.80 · F1 >= 0.70 · FP/h <= 10 · gap val-test <= 0.08",
        0.3, 6.35, 12.7, 0.38, size=10, italic=True, color=AM,
        align=PP_ALIGN.CENTER)

    # barra visual AUC-PR
    box(s, 0.3, 6.82, 12.7, 0.5, fill=BR, border=CZ2, border_w=1)
    txt(s, "AUC-PR:", 0.4, 6.87, 1.0, 0.35, size=9, color=CZ2)
    modelos_bar = [("Tiny CNN atual", 0.8775, VE),
                   ("Random Forest", 0.8127, AM),
                   ("STA/LTA", 0.1662, VR)]
    for label, val, c in modelos_bar:
        pass  # barra nao cabe — skip


def slide_quality_gate(prs, n, tot):
    s = slide_branco(prs, CZ)
    header(s, "Quality Gate", "Criterios automaticos para promocao a producao")
    numero_slide(s, n, tot)

    # fluxo
    etapas_qg = ["candidate_manifest.json", "promote_model.py",
                 "Verifica criterios", "production_manifest.json"]
    for i, e in enumerate(etapas_qg):
        cor = [AM, AM_, VE, VE][i]
        chip(s, e, 1.0 + i*3.05, 1.45, w=2.7, h=0.58, fill=cor, size=10)
        if i < len(etapas_qg)-1:
            txt(s, "→", 3.72+i*3.05, 1.6, 0.35, 0.3,
                size=14, bold=True, color=AM, align=PP_ALIGN.CENTER)

    # criterios
    box(s, 0.3, 2.2, 6.2, 4.9, fill=BR, border=AZ, border_w=1)
    txt(s, "Criterios Configurados", 0.4, 2.23, 6.0, 0.38,
        size=12, bold=True, color=AZ)

    criterios = [
        ("AUC-PR minima",         ">= 0.80",  "config.yaml"),
        ("F1 minimo",             ">= 0.70",  "config.yaml"),
        ("FP por hora",           "<= 10.0",  "config.yaml"),
        ("Gap val-test AUC-PR",   "<= 0.08",  "anti-overfitting"),
        ("Tamanho do modelo",     "<= 300 KB","viabilidade edge"),
    ]
    for i,(k,v,obs) in enumerate(criterios):
        y = 2.75 + i*0.75
        box(s, 0.4, y, 3.5, 0.62, fill=AZ)
        txt(s, k, 0.45, y+0.07, 3.4, 0.5, size=10, bold=True, color=BR)
        box(s, 3.95, y, 2.4, 0.62, fill=AC)
        txt(s, v, 3.97, y+0.07, 2.3, 0.5,
            size=12, bold=True, color=AZ, align=PP_ALIGN.CENTER)

    # resultado
    box(s, 6.8, 2.2, 6.2, 4.9, fill=BR, border=VE, border_w=2)
    box(s, 6.8, 2.2, 6.2, 0.45, fill=VE)
    txt(s, "Resultado — Tiny CNN APROVADO", 6.9, 2.22, 6.0, 0.4,
        size=12, bold=True, color=BR)

    resultados = [
        ("AUC-PR", "0.8775", ">= 0.80", True),
        ("F1", "0.8109", ">= 0.70", True),
        ("FP/h", "6.75", "<= 10.0", True),
        ("Gap val-test", "0.018", "<= 0.08", True),
        ("Tamanho", "25.7 KB", "<= 300 KB", True),
    ]
    for i,(k,v,regra,ok) in enumerate(resultados):
        y = 2.75 + i*0.75
        box(s, 6.9, y, 2.5, 0.62, fill=CZ)
        txt(s, k, 6.95, y+0.08, 2.4, 0.48, size=11, color=PT, bold=True)
        box(s, 9.45, y, 1.7, 0.62, fill=VA if ok else RGBColor(0xF9,0xEB,0xEA))
        txt(s, v, 9.47, y+0.08, 1.65, 0.48,
            size=13, bold=True, color=VE if ok else VR,
            align=PP_ALIGN.CENTER)
        txt(s, "✅ "+regra if ok else "❌ "+regra,
            11.2, y+0.1, 1.7, 0.45, size=9,
            color=VE if ok else VR)

    txt(s, "Todos os 5 criterios aprovados — modelo promovido a producao",
        0.3, 7.2, 12.7, 0.28, size=11, bold=True,
        color=VE, align=PP_ALIGN.CENTER)


def slide_export_tflite(prs, n, tot):
    s = slide_branco(prs, CZ)
    header(s, "Exportacao TFLite", "Tres formatos de quantizacao para diferentes cenarios")
    numero_slide(s, n, tot)

    # pipeline export
    steps_exp = [("tiny_cnn.keras\n(modelo Keras)", AM),
                 ("export_tflite.py\n(conversao)", AM_),
                 ("Representative\nDataset (int8)", AC),
                 ("Artefatos\nedge/", VE)]
    for i,(lbl,c) in enumerate(steps_exp):
        chip(s, lbl, 0.35+i*3.25, 1.45, w=3.0, h=0.72, fill=c, size=10)
        if i < len(steps_exp)-1:
            txt(s, "→", 3.37+i*3.25, 1.65, 0.25, 0.3,
                size=14, bold=True, color=AM, align=PP_ALIGN.CENTER)

    # tres formatos
    formatos = [
        ("float32.tflite", "66.2 KB", "Pesos em 32 bits\nNenhuma quantizacao",
         "Firmware atual\n(validado no ESP32)", AM,
         "Precisao maxima\nMaior consumo RAM"),
        ("float16.tflite", "37.9 KB", "Pesos em 16 bits\nQuantizacao de pesos",
         "Alternativa compacta", AC,
         "Bom equilibrio\ntamanho/precisao"),
        ("int8.tflite", "25.7 KB", "Pesos e ativacoes em 8 bits\nPTQ com rep. dataset",
         "Alvo final\n(REDUCE_MAX pendente)", VE,
         "Menor tamanho\nMaior velocidade"),
    ]
    for i,(nome,tam,desc,uso,c,nota) in enumerate(formatos):
        x = 0.3 + i*4.35
        box(s, x, 2.38, 4.1, 4.55, fill=BR, border=c, border_w=2)
        box(s, x, 2.38, 4.1, 0.52, fill=c)
        txt(s, nome, x+0.1, 2.4, 3.9, 0.48,
            size=13, bold=True, color=BR, align=PP_ALIGN.CENTER)
        box(s, x+0.8, 3.0, 2.5, 0.55, fill=c)
        txt(s, tam, x+0.8, 3.0, 2.5, 0.55,
            size=20, bold=True, color=BR, align=PP_ALIGN.CENTER)
        txt(s, desc, x+0.1, 3.65, 3.9, 0.65,
            size=10, color=PT, align=PP_ALIGN.CENTER)
        box(s, x+0.1, 4.4, 3.9, 0.52, fill=CZ)
        txt(s, uso, x+0.15, 4.42, 3.8, 0.48,
            size=10, italic=True, color=AM, align=PP_ALIGN.CENTER)
        txt(s, nota, x+0.1, 5.02, 3.9, 0.6,
            size=9, color=CZ2, align=PP_ALIGN.CENTER)

    txt(s,
        "Header C/C++ (.h) gerado automaticamente — incluso diretamente no firmware PlatformIO",
        0.3, 7.1, 12.7, 0.35, size=11, italic=True,
        color=AM, align=PP_ALIGN.CENTER)


def slide_drift(prs, n, tot):
    s = slide_branco(prs, CZ)
    header(s, "Drift Detection", "Monitoramento de distribuicao e ciclo de retreino")
    numero_slide(s, n, tot)

    # diagrama de fluxo
    etapas_d = ["build_drift\n_reference", "check_data\n_drift",
                "retrain\n_policy", "drift_to_ota\n_decision"]
    for i,e in enumerate(etapas_d):
        chip(s, e, 0.35+i*3.25, 1.45, w=3.0, h=0.65, fill=AM, size=9)
        if i < len(etapas_d)-1:
            txt(s, "→", 3.37+i*3.25, 1.58, 0.25, 0.3,
                size=14, bold=True, color=AM, align=PP_ALIGN.CENTER)

    # metricas
    metricas_d = [
        ("z-shift (max)", "0.0176",
         "Deslocamento de media\nem desvios padrao",
         "Baixo — medias\npouco alteradas", CZ2),
        ("PSI (max)", "0.3463",
         "Population Stability Index\nmudanca de distribuicao",
         "ALTO\n> 0.2 = significativo", VR),
        ("KS p-value (min)", "0.000033",
         "Teste Kolmogorov-Smirnov\ndiferenca estatistica",
         "Muito significativo\np << 0.05", VR),
    ]
    for i,(k,v,desc,interp,c) in enumerate(metricas_d):
        x = 0.3 + i*4.35
        box(s, x, 2.3, 4.1, 3.3, fill=BR, border=c, border_w=2)
        box(s, x, 2.3, 4.1, 0.45, fill=c)
        txt(s, k, x+0.1, 2.32, 3.9, 0.4,
            size=11, bold=True, color=BR, align=PP_ALIGN.CENTER)
        txt(s, v, x+0.5, 2.85, 3.1, 0.7,
            size=26, bold=True, color=c if c!=CZ2 else PT,
            align=PP_ALIGN.CENTER)
        txt(s, desc, x+0.1, 3.6, 3.9, 0.6,
            size=9, color=PT, align=PP_ALIGN.CENTER)
        box(s, x+0.1, 4.25, 3.9, 0.6, fill=CZ)
        txt(s, interp, x+0.15, 4.27, 3.8, 0.55,
            size=10, bold=True,
            color=VR if c==VR else CZ2,
            align=PP_ALIGN.CENTER)

    # decisao
    box(s, 0.3, 5.8, 5.9, 1.05, fill=BR, border=AM_, border_w=2)
    box(s, 0.3, 5.8, 5.9, 0.42, fill=AM_)
    txt(s, "Politica", 0.4, 5.82, 5.7, 0.38, size=11, bold=True, color=BR)
    txt(s, "retrain_recommended", 0.4, 6.3, 5.7, 0.45,
        size=14, bold=True, color=AM_, align=PP_ALIGN.CENTER)

    box(s, 6.55, 5.8, 6.5, 1.05, fill=BR, border=VE, border_w=2)
    box(s, 6.55, 5.8, 6.5, 0.42, fill=VE)
    txt(s, "Decisao OTA", 6.65, 5.82, 6.3, 0.38,
        size=11, bold=True, color=BR)
    txt(s, "build_and_publish_ota", 6.65, 6.3, 6.3, 0.45,
        size=14, bold=True, color=VE, align=PP_ALIGN.CENTER)

    txt(s, "OTA so liberado quando candidato novo passa o quality gate",
        0.3, 7.0, 12.7, 0.35, size=10, italic=True,
        color=CZ2, align=PP_ALIGN.CENTER)


def slide_ota_simulado(prs, n, tot):
    s = slide_branco(prs, CZ)
    header(s, "Fluxo OTA Simulado (MLOps)", "Manifesto · Pacote · Assinatura HMAC-SHA256 · Rollback")
    numero_slide(s, n, tot)

    # etapas em duas colunas
    etapas = [
        ("build_ota_manifest", "Gera ota_manifest.json a partir\ndo production_manifest", AM),
        ("build_ota_package", "Empacota artifact.tflite\n+ calcula SHA-256", AM_),
        ("validate_ota_package", "Verifica integridade:\nSHA-256 + assinatura HMAC", VE),
        ("publish_local_release", "Cria releases/latest.json\n(simula repositorio OTA)", AC),
        ("simulate_device_update_check", "Dispositivo consulta latest.json\ne verifica compatibilidade", AM),
        ("simulate_apply_update", "Instalacao simulada\ncom log completo", VE),
        ("simulate_rollback", "Reversao automatica\nem caso de falha", VR),
    ]
    for i,(cmd,desc,c) in enumerate(etapas):
        col = i % 2; row = i // 2
        x = 0.3 + col*6.65
        y = 1.48 + row*1.42
        box(s, x, y, 6.3, 1.28, fill=BR, border=c, border_w=1)
        box(s, x, y, 6.3, 0.45, fill=c)
        txt(s, f"python -m src.ota.{cmd}",
            x+0.1, y+0.05, 6.1, 0.38, size=9, bold=True, color=BR)
        txt(s, desc, x+0.1, y+0.55, 6.1, 0.65, size=10, color=PT)

    # badges de resultado
    badges = [("SHA-256 OK", VE), ("HMAC OK", VE),
              ("Rollback OK", AM_), ("Versao: seismic_edge_v1_tiny_cnn_20260614", AM)]
    for i,(b,c) in enumerate(badges):
        chip(s, b, 0.3+i*3.25, 7.05, w=3.05, h=0.38,
             fill=c, size=9)


def slide_ota_real(prs, n, tot):
    s = slide_branco(prs, CZ)
    header(s, "OTA Real no ESP32 via HTTP", "Firmware consulta servidor, baixa modelo, valida e carrega")
    numero_slide(s, n, tot)

    # lado esquerdo: servidor
    box(s, 0.3, 1.42, 5.5, 5.75, fill=BR, border=AM, border_w=2)
    box(s, 0.3, 1.42, 5.5, 0.48, fill=AM)
    txt(s, "Servidor FastAPI (:8000)", 0.4, 1.44, 5.3, 0.43,
        size=12, bold=True, color=BR)

    endpoints = [
        ("GET /ota/latest", "Retorna latest.json com versao e SHA-256"),
        ("GET /ota/artifact", "Faz download do .tflite (streaming)"),
        ("POST /ota/report", "Recebe resultado: success/failed/skipped"),
        ("GET /metrics", "Metricas Prometheus"),
        ("GET /devices", "Registry de dispositivos"),
    ]
    for i,(ep,desc) in enumerate(endpoints):
        y = 2.05 + i*0.95
        box(s, 0.4, y, 2.3, 0.75, fill=AM)
        txt(s, ep, 0.42, y+0.1, 2.28, 0.58, size=8, bold=True,
            color=BR, align=PP_ALIGN.CENTER)
        txt(s, desc, 2.78, y+0.12, 2.9, 0.55, size=9, color=PT)

    # seta central
    for y in [3.0, 4.0, 5.0]:
        txt(s, "⇄", 5.95, y, 0.7, 0.5,
            size=22, bold=True, color=AM_, align=PP_ALIGN.CENTER)
    txt(s, "HTTP\nWiFi", 5.85, 4.0, 0.8, 0.6,
        size=9, color=CZ2, align=PP_ALIGN.CENTER)

    # lado direito: esp32 fluxo
    box(s, 7.0, 1.42, 6.0, 5.75, fill=BR, border=VE, border_w=2)
    box(s, 7.0, 1.42, 6.0, 0.48, fill=VE)
    txt(s, "ESP32 — ota_http.cpp", 7.1, 1.44, 5.8, 0.43,
        size=12, bold=True, color=BR)

    fluxo = [
        ("SPIFFS.begin()", AC, "Monta filesystem na flash"),
        ("wifi_connect()", AM, "Conecta ao AP (15s timeout)"),
        ("GET /ota/latest", AM_, "Compara versoes"),
        ("GET /ota/artifact", AM_, "Stream → SPIFFS (512B/vez)"),
        ("sha256_of_file()", VE, "mbedtls — valida integridade"),
        ("ota_load_into_ram()", VE, "heap_caps_aligned_alloc(16B)"),
        ("tflite::GetModel()", AZ, "Ponteiro SPIFFS ou builtin"),
    ]
    for i,(cmd,c,desc) in enumerate(fluxo):
        y = 2.05 + i*0.73
        box(s, 7.1, y, 2.5, 0.6, fill=c)
        txt(s, cmd, 7.12, y+0.08, 2.48, 0.46,
            size=8, bold=True, color=BR, align=PP_ALIGN.CENTER)
        txt(s, desc, 9.68, y+0.1, 3.2, 0.45, size=9, color=PT)

    txt(s,
        "Fallback: se WiFi falhar ou SHA-256 divergir → usa modelo compilado no header (.h)",
        0.3, 7.25, 12.7, 0.3, size=10, italic=True,
        color=CZ2, align=PP_ALIGN.CENTER)


def slide_memoria(prs, n, tot):
    s = slide_branco(prs, CZ)
    header(s, "Mapa de Memoria — ESP32", "Flash e RAM com firmware TFLite Micro + OTA")
    numero_slide(s, n, tot)

    # flash
    box(s, 0.3, 1.42, 5.8, 5.8, fill=BR, border=AM, border_w=2)
    box(s, 0.3, 1.42, 5.8, 0.48, fill=AM)
    txt(s, "FLASH — 4 MB (huge_app.csv)", 0.4, 1.44, 5.6, 0.43,
        size=12, bold=True, color=BR)

    flash_segs = [
        ("Bootloader", "8 KB", 0x04, CZ2),
        ("NVS (WiFi config)", "20 KB", 0.12, CZ2),
        ("OTA data", "8 KB", 0.04, CZ2),
        ("app0 — Firmware", "3 MB  (usado: 638 KB → 20.3%)", 0.45, AM),
        ("SPIFFS — dados", "960 KB  (ota_model.tflite: 66 KB)", 0.2, AC),
    ]
    y = 2.05
    for nome,tam,frac,c in flash_segs:
        h_seg = max(frac*4.6, 0.42)
        box(s, 0.4, y, 5.5, h_seg, fill=c, border=BR, border_w=1)
        txt(s, nome, 0.5, y+0.05, 3.2, h_seg-0.1,
            size=9, bold=True, color=BR if c!=CZ2 else PT)
        txt(s, tam, 3.75, y+0.05, 2.0, h_seg-0.1,
            size=8, color=BR if c!=CZ2 else PT, align=PP_ALIGN.RIGHT)
        y += h_seg

    # ram
    box(s, 6.7, 1.42, 6.3, 5.8, fill=BR, border=VE, border_w=2)
    box(s, 6.7, 1.42, 6.3, 0.48, fill=VE)
    txt(s, "RAM — 320 KB", 6.8, 1.44, 6.1, 0.43,
        size=12, bold=True, color=BR)

    ram_segs = [
        ("Sistema + Arduino + codigo", "~50 KB", 0.16, CZ2),
        ("tensor_arena (TFLite)", "100 KB  (alocado pelo fw)", 0.31, AM),
        ("s_model_buf (modelo OTA)", "66 KB  (heap_caps 16B)", 0.21, AC),
        ("Disponivel", "~104 KB  livres", 0.32, VA),
    ]
    y = 2.05
    for nome,tam,frac,c in ram_segs:
        h_seg = max(frac*4.6, 0.42)
        box(s, 6.8, y, 6.1, h_seg, fill=c, border=BR, border_w=1)
        txt(s, nome, 6.9, y+0.05, 3.8, h_seg-0.1,
            size=9, bold=True, color=BR if c not in (CZ2,VA) else PT)
        txt(s, tam, 10.6, y+0.05, 2.1, h_seg-0.1,
            size=8, color=BR if c not in (CZ2,VA) else PT,
            align=PP_ALIGN.RIGHT)
        y += h_seg

    txt(s, "RAM total usada (firmware + OTA): ~216 KB / 320 KB = 67.5%  ·  Flash: 638 KB / 3072 KB = 20.3%",
        0.3, 7.3, 12.7, 0.3, size=10, italic=True,
        color=AM, align=PP_ALIGN.CENTER)


def slide_validacao_hw(prs, n, tot):
    s = slide_branco(prs, CZ)
    header(s, "Validacao no Hardware Real", "ESP32-D0WD-V3 · TFLite Micro · Inferencia embarcada")
    numero_slide(s, n, tot)

    # checklist
    box(s, 0.3, 1.42, 7.8, 5.8, fill=BR, border=VE, border_w=2)
    box(s, 0.3, 1.42, 7.8, 0.48, fill=VE)
    txt(s, "Status de Validacao", 0.4, 1.44, 7.6, 0.43,
        size=12, bold=True, color=BR)

    checks = [
        ("OK", "ESP32 detectado no WSL via /dev/ttyUSB0", VE),
        ("OK", "Chip identificado: ESP32-D0WD-V3", VE),
        ("OK", "Firmware compilado: RAM 31% / Flash 20%", VE),
        ("OK", "Upload via esptool bem-sucedido", VE),
        ("OK", "Modelo float32 carregado e invocado", VE),
        ("OK", "Saida CSV no serial (score, latencia, CPU%)", VE),
        ("OK", "OTA HTTP: download + SHA-256 + SPIFFS", VE),
        ("⚠", "Inferencia int8: REDUCE_MAX (correcao planejada)", AM_),
    ]
    for i,(status,desc,c) in enumerate(checks):
        y = 2.05 + i*0.65
        box(s, 0.4, y, 0.6, 0.55, fill=c)
        txt(s, status, 0.4, y+0.05, 0.6, 0.48,
            size=11, bold=True, color=BR, align=PP_ALIGN.CENTER)
        txt(s, desc, 1.1, y+0.08, 6.8, 0.45, size=11, color=PT)

    # serial output
    box(s, 8.4, 1.42, 4.9, 5.8, fill=RGBColor(0x1E,0x1E,0x1E))
    box(s, 8.4, 1.42, 4.9, 0.48, fill=PT)
    txt(s, "Saida Serial (monitor)", 8.5, 1.44, 4.7, 0.43,
        size=11, bold=True, color=CZ)

    serial_lines = [
        "[WiFi] IP: 192.168.1.107",
        "[OTA] Nova versao disponivel",
        "[OTA] Baixados: 67768 bytes",
        "[OTA] SHA-256 OK.",
        "Modelo: PIPELINE_TINY_CNN_FLOAT32",
        "Threshold: 0.72419429",
        "run,score,pred,correct,...",
        "0,0.8234,1,1,...",
        "1,0.1023,0,1,...",
        "# summary,inference_ms",
        "# avg: 8.43 ms",
        "# summary,accuracy,0.85",
    ]
    for i,line in enumerate(serial_lines):
        cor = VE if line.startswith("[OTA]") or "OK" in line else \
              AM_ if "summary" in line or "Threshold" in line else \
              AC if line.startswith("Modelo") else CZ
        txt(s, line, 8.5, 2.05+i*0.43, 4.7, 0.4,
            size=8, color=cor)


def slide_resultados(prs, n, tot):
    s = slide_branco(prs, CZ)
    header(s, "Resultados Finais", "Tiny CNN — 15.377 parametros — dataset sismico")
    numero_slide(s, n, tot)

    # metricas grandes
    metricas_g = [
        ("AUC-PR", "0.8775", "teste", VE),
        ("AUC-ROC", "0.9576", "teste", AM),
        ("F1", "0.8109", "teste", AM_),
        ("FP/h", "6.75", "teste", AC),
    ]
    for i,(k,v,split,c) in enumerate(metricas_g):
        x = 0.3 + i*3.25
        box(s, x, 1.42, 3.05, 2.1, fill=BR, border=c, border_w=2)
        box(s, x, 1.42, 3.05, 0.48, fill=c)
        txt(s, k, x+0.1, 1.44, 2.9, 0.43,
            size=13, bold=True, color=BR, align=PP_ALIGN.CENTER)
        txt(s, v, x+0.1, 1.98, 2.9, 0.9,
            size=32, bold=True, color=c, align=PP_ALIGN.CENTER)
        txt(s, f"conjunto de {split}", x+0.1, 2.95, 2.9, 0.35,
            size=9, color=CZ2, align=PP_ALIGN.CENTER)

    # tabela val vs teste
    box(s, 0.3, 3.72, 6.0, 3.2, fill=BR, border=AM, border_w=1)
    txt(s, "Validacao vs Teste", 0.4, 3.75, 5.8, 0.38,
        size=12, bold=True, color=AM)

    rows = [("Metrica","Validacao","Teste"),
            ("AUC-PR","0.8598","0.8775"),
            ("AUC-ROC","0.9546","0.9576"),
            ("F1","0.7908","0.8109"),
            ("Precision","0.8343","0.8431"),
            ("Recall","0.7517","0.7810")]
    alts2 = [2.5, 1.7, 1.8]
    for i,row in enumerate(rows):
        bg = AZ if i==0 else (VA if i%2==0 else BR)
        fg = BR if i==0 else PT
        linha_tabela(s, None, row, 4.2+i*0.42, alts2,
                     cor_bg=bg, cor_txt=fg,
                     size=10, bold=(i==0))

    # confusion matrix simplificada
    box(s, 6.65, 3.72, 6.35, 3.2, fill=BR, border=AM, border_w=1)
    txt(s, "Confusion Matrix (teste)", 6.75, 3.75, 6.15, 0.38,
        size=12, bold=True, color=AM)

    # matriz 2x2
    labels = [["","Pred 0","Pred 1"],
              ["Real 0","TN: 22.016","FP: 484"],
              ["Real 1","FN: 729","TP: 2.600"]]
    cores_m = [[PT,CZ2,CZ2],[CZ2,VA,RGBColor(0xF9,0xEB,0xEA)],
               [CZ2,RGBColor(0xF9,0xEB,0xEA),VA]]
    for r,row in enumerate(labels):
        for c,cel in enumerate(row):
            bx = 6.75 + c*1.95
            by = 4.25 + r*0.8
            box(s, bx, by, 1.88, 0.75,
                fill=cores_m[r][c], border=CZ2, border_w=0.5)
            txt(s, cel, bx+0.05, by+0.1, 1.78, 0.55,
                size=10, bold=(r==0 or c==0),
                color=VE if "TP" in cel or "TN" in cel else
                      VR if "FP" in cel or "FN" in cel else PT,
                align=PP_ALIGN.CENTER)

    txt(s, "Threshold 0.724 escolhido no val · aplicado diretamente no teste",
        0.3, 7.0, 12.7, 0.38, size=10, italic=True,
        color=CZ2, align=PP_ALIGN.CENTER)


def slide_multivariado(prs, n, tot):
    s = slide_branco(prs, CZ)
    header(s, "Proxima Direcao: Series Temporais Multivariadas",
           "Nova branch em desenvolvimento")
    numero_slide(s, n, tot)

    # comparativo univariado vs multivariado
    box(s, 0.3, 1.42, 6.0, 5.8, fill=BR, border=AM, border_w=2)
    box(s, 0.3, 1.42, 6.0, 0.48, fill=AM)
    txt(s, "Atual — Univariado", 0.4, 1.44, 5.8, 0.43,
        size=13, bold=True, color=BR)

    uni = [("Dataset shape","X.shape = (n, 800)"),
           ("Canais","1 canal (ex: aceleracao Z)"),
           ("Preprocessamento","zscore global por janela"),
           ("Modelo","Conv1D 1D simples"),
           ("Tensor input","shape = (800, 1)"),
           ("RAM modelo","66 KB (float32)")]
    for i,(k,v) in enumerate(uni):
        y = 2.07+i*0.82
        box(s, 0.4, y, 2.3, 0.68, fill=AZ)
        txt(s, k, 0.43, y+0.09, 2.25, 0.52,
            size=10, bold=True, color=BR, align=PP_ALIGN.CENTER)
        txt(s, v, 2.78, y+0.12, 3.3, 0.5, size=10, color=PT)

    box(s, 6.7, 1.42, 6.3, 5.8, fill=BR, border=AM_, border_w=2)
    box(s, 6.7, 1.42, 6.3, 0.48, fill=AM_)
    txt(s, "Novo — Multivariado", 6.8, 1.44, 6.1, 0.43,
        size=13, bold=True, color=BR)

    multi = [("Dataset shape","X.shape = (n, 800, C)"),
             ("Canais","C canais (X, Y, Z acelerometro...)"),
             ("Preprocessamento","zscore por canal ou global"),
             ("Modelo","Conv2D ou TCN multicanal"),
             ("Tensor input","shape = (800, C)"),
             ("RAM modelo","maior — validar ESP32")]
    for i,(k,v) in enumerate(multi):
        y = 2.07+i*0.82
        box(s, 6.8, y, 2.3, 0.68, fill=AM_)
        txt(s, k, 6.83, y+0.09, 2.25, 0.52,
            size=10, bold=True, color=BR, align=PP_ALIGN.CENTER)
        txt(s, v, 9.18, y+0.12, 3.6, 0.5, size=10, color=PT)

    txt(s,
        "Pipeline generico foi projetado para esta extensao — adapter de dominio isola a diferenca de contrato",
        0.3, 7.2, 12.7, 0.3, size=11, italic=True,
        color=AM, align=PP_ALIGN.CENTER)


def slide_proximos_passos(prs, n, tot):
    s = slide_branco(prs, CZ)
    header(s, "Proximos Passos", "Prioridades tecnicas e evolucao do sistema")
    numero_slide(s, n, tot)

    passos = [
        ("ALTA", "Correcao int8",
         "Retreinar com head_pooling=avg (ja configurado no YAML) e reexportar.\nTensor arena int8 cabe em 100 KB — 3x menos que float32.",
         VR),
        ("ALTA", "Validacao int8 real no ESP32",
         "Confirmar inferencia int8 no hardware apos correcao.\nMedir latencia e comparar com float32.",
         VR),
        ("MEDIA", "Reset automatico pos-OTA",
         "Chamar ESP.restart() apos OTA bem-sucedida.\nGarante ambiente de execucao limpo.",
         AM_),
        ("MEDIA", "Series Temporais Multivariadas",
         "Expandir pipeline para multiplos canais (branch em desenvolvimento).\nAdapter de dominio ja isola a diferenca de contrato.",
         AM_),
        ("MEDIA", "Rollback automatico de firmware",
         "Se AllocateTensors() falhar com modelo OTA,\ndeletar SPIFFS e reverter para builtin automaticamente.",
         AM_),
        ("BAIXA", "Assinatura RSA/ECDSA no firmware",
         "Substituir HMAC-SHA256 por criptografia assimetrica.\nChave publica no dispositivo — sem segredo compartilhado.",
         AC),
        ("FUTURA", "Artigo de sistema TinyML/MLOps",
         "Documentar pipeline, resultados e comparativo de modelos\nem formato de publicacao cientifica.",
         CZ2),
    ]
    for i,(_,titulo,desc,c) in enumerate(passos):
        col = i % 2; row = i // 2
        if i == 6: col=0; row=3
        x = 0.3 + col*6.65
        y = 1.48 + row*1.42
        if i < 6:
            box(s, x, y, 6.3, 1.28, fill=BR, border=c, border_w=2)
            box(s, x, y, 6.3, 0.45, fill=c)
            txt(s, titulo, x+0.1, y+0.06, 6.1, 0.38,
                size=11, bold=True, color=BR)
            txt(s, desc, x+0.1, y+0.55, 6.1, 0.65, size=9, color=PT)

    # ultimo item no centro
    x = 0.3 + 6.65
    y = 1.48 + 3*1.42
    box(s, x, y, 6.3, 1.28, fill=BR, border=CZ2, border_w=2)
    box(s, x, y, 6.3, 0.45, fill=CZ2)
    txt(s, "Artigo de sistema TinyML/MLOps",
        x+0.1, y+0.06, 6.1, 0.38, size=11, bold=True, color=BR)
    txt(s, "Documentar pipeline, resultados e comparativo de modelos em publicacao cientifica.",
        x+0.1, y+0.55, 6.1, 0.65, size=9, color=PT)


def slide_conclusao(prs, n, tot):
    s = slide_branco(prs, AZ)
    box(s, 0, 0, 0.2, 7.5, fill=AM_)
    box(s, 0.2, 0, 13.13, 7.5, fill=AZ)
    box(s, 0.2, 5.55, 13.13, 1.95, fill=AM)

    txt(s, "Conclusao", 0.5, 0.3, 12.3, 0.65,
        size=32, bold=True, color=BR)
    box(s, 0.5, 0.95, 7.0, 0.06, fill=AM_)

    conquistas = [
        ("Pipeline MLOps completo",
         "Do dado bruto ao dispositivo — rastreavel, automatizado e reproduzivel"),
        ("Modelo tiny_cnn aprovado",
         "AUC-PR 0.877 · F1 0.811 · FP/h 6.75 · quality gate aprovado"),
        ("Validacao fisica real",
         "Build, flash e inferencia float32 no ESP32-D0WD-V3"),
        ("OTA via HTTP implementado",
         "WiFi + download streaming + SHA-256 + SPIFFS + fallback builtin"),
        ("Drift detection integrado",
         "PSI 0.346, KS p<0.001 → retrain → OTA → ciclo fechado"),
    ]
    for i,(titulo,desc) in enumerate(conquistas):
        y = 1.18 + i*0.82
        box(s, 0.5, y, 0.5, 0.6, fill=VE)
        txt(s, "✓", 0.5, y+0.08, 0.5, 0.48,
            size=14, bold=True, color=BR, align=PP_ALIGN.CENTER)
        txt(s, titulo, 1.1, y+0.02, 4.2, 0.3,
            size=12, bold=True, color=BR)
        txt(s, desc, 1.1, y+0.32, 11.2, 0.28,
            size=10, color=RGBColor(0xA8,0xC8,0xE8))

    txt(s,
        "Proximo: series temporais multivariadas  ·  int8 corrigido  ·  OTA completo",
        0.5, 5.75, 12.3, 0.42, size=13, color=BR, align=PP_ALIGN.CENTER)
    txt(s, "github.com/alvarosamp/TCC  ·  2026",
        0.5, 6.25, 12.3, 0.38, size=12,
        color=RGBColor(0x90,0xB8,0xD8), align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slides_fn = [
        slide_capa,
        slide_problema,
        slide_arquitetura,
        slide_dataset,
        slide_modelos_arch,
        slide_treinamento,
        slide_comparativo,
        slide_quality_gate,
        slide_export_tflite,
        slide_drift,
        slide_ota_simulado,
        slide_ota_real,
        slide_memoria,
        slide_validacao_hw,
        slide_resultados,
        slide_multivariado,
        slide_proximos_passos,
        slide_conclusao,
    ]
    total = len(slides_fn)

    slide_capa(prs)
    for i, fn in enumerate(slides_fn[1:], start=2):
        fn(prs, i, total)

    saida = os.path.abspath(
        os.path.join(os.path.dirname(__file__), "..", "docs", "slides_tcc.pptx")
    )
    prs.save(saida)
    print(f"Slides salvos em: {saida}")
    print(f"Total de slides : {total}")


if __name__ == "__main__":
    main()
